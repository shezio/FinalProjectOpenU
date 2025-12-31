"""
Services for dashboard video generation and PPT export
"""
import os
import tempfile
import uuid
from pathlib import Path
from datetime import timedelta
from django.utils import timezone
from django.db.models import Count, Q
import threading
import time

from .models import Children, Tutors, Tutorships, Feedback, Staff


def generate_dashboard_data(timeframe):
    """Generate dashboard data based on timeframe"""
    
    now = timezone.now()
    
    # Calculate date range
    if timeframe == 'שבוע אחרון':
        start_date = now - timedelta(days=7)
    elif timeframe == 'חודש אחרון' or timeframe == '3 חודשים אחרונים':
        start_date = now - timedelta(days=30)
    elif timeframe == 'שנה אחרונה':
        start_date = now - timedelta(days=365)
    else:  # 'כל הזמן'
        start_date = None
    
    # Collect data
    data = {
        'total_families': Children.objects.count(),
        'waiting_families': Children.objects.filter(
            tutorships__isnull=True
        ).count(),
        'active_tutorships': Tutorships.objects.count(),
        'pending_tutors': Tutors.objects.filter(tutorship_status='אין_חניך').count(),
        'staff_count': Staff.objects.filter(registration_approved=True).count(),
        'new_families': Children.objects.filter(registrationdate__gte=start_date).count() if start_date else 0,
        'timeframe': timeframe
    }
    
    return data


def generate_ai_video(video_id, dashboard_data, timeframe, duration, pages, style):
    """
    Generate AI marketing video with Hebrew narration using gTTS and MoviePy
    Creates professional slides with system data and Hebrew voiceover
    """
    try:
        from gtts import gTTS
        from PIL import Image, ImageDraw, ImageFont
        from moviepy import ImageSequenceClip, AudioFileClip, concatenate_videoclips
    except ImportError as e:
        print(f"❌ Import error: {e}")
        raise Exception(f"Missing dependency: {e}")

    # Create temp directory
    temp_dir = Path(tempfile.gettempdir()) / 'childsmile_videos'
    temp_dir.mkdir(exist_ok=True)
    
    video_path = temp_dir / f"{video_id}.mp4"
    audio_path = temp_dir / f"{video_id}_audio.mp3"
    
    try:
        # Generate marketing script in Hebrew
        script = generate_marketing_script(dashboard_data, timeframe, style)
        print(f"✅ Generated script ({len(script)} chars)")
        
        # Create TTS audio
        print("⏳ Creating TTS audio...")
        tts = gTTS(text=script, lang='iw', slow=False)  # 'iw' is the language code for Hebrew
        tts.save(str(audio_path))
        print(f"✅ Audio saved: {audio_path}")
        
        # Create professional slides
        print("⏳ Creating slides...")
        slides_images = create_professional_slides(dashboard_data, timeframe)
        print(f"✅ Created {len(slides_images)} slides")
        
        # Save slides as temporary image files
        slide_paths = []
        for i, slide_img in enumerate(slides_images):
            slide_path = temp_dir / f"{video_id}_slide_{i}.png"
            slide_img.save(str(slide_path))
            slide_paths.append(str(slide_path))
        print(f"✅ Saved {len(slide_paths)} slide images")
        
        # Load audio first to get duration
        print("⏳ Loading audio...")
        audio_clip = AudioFileClip(str(audio_path))
        audio_duration = audio_clip.duration
        print(f"✅ Audio duration: {audio_duration}s")
        
        # Calculate slide durations intelligently
        # First and last slides: 2 seconds each
        # Middle slides: share the remaining time evenly
        num_slides = len(slide_paths)
        first_slide_duration = 2.0
        last_slide_duration = 2.0
        
        if num_slides <= 2:
            # If 1-2 slides, divide the audio duration evenly
            slide_durations = [audio_duration / num_slides] * num_slides
        else:
            # Reserve 2 seconds for first and last slides
            remaining_duration = audio_duration - first_slide_duration - last_slide_duration
            middle_slides_count = num_slides - 2
            middle_slide_duration = remaining_duration / middle_slides_count
            
            # Build duration list
            slide_durations = [first_slide_duration]
            slide_durations.extend([middle_slide_duration] * middle_slides_count)
            slide_durations.append(last_slide_duration)
        
        print(f"⏳ Slide durations: {[f'{d:.1f}s' for d in slide_durations]}")
        
        # Create video from slides with calculated durations
        print("⏳ Creating video clips...")
        slide_clips = []
        for slide_path, duration in zip(slide_paths, slide_durations):
            clip = ImageSequenceClip([slide_path], fps=1)
            clip = clip.with_duration(duration)
            slide_clips.append(clip)
        print(f"✅ Created {len(slide_clips)} clip objects")
        
        # Concatenate all slides
        print("⏳ Concatenating clips...")
        video_clip = concatenate_videoclips(slide_clips, method="compose")
        print(f"✅ Video clip duration: {video_clip.duration}s")
        
        # Adjust video duration to match audio
        video_duration = video_clip.duration
        audio_duration = audio_clip.duration
        
        if video_duration < audio_duration:
            # Video is shorter - extend it by repeating the last frame
            print(f"⏳ Extending video from {video_duration}s to {audio_duration}s...")
            video_clip = video_clip.with_duration(audio_duration)
        elif video_duration > audio_duration:
            # Video is longer - trim it
            print(f"⏳ Trimming video from {video_duration}s to {audio_duration}s...")
            video_clip = video_clip.with_duration(audio_duration)
        
        print(f"✅ Synchronized video duration: {video_clip.duration}s")
        
        # Set audio on video using with_audio (MoviePy 2.x API)
        final_clip = video_clip.with_audio(audio_clip)
        print("✅ Audio attached to video")
        
        # Write video file
        print(f"⏳ Writing MP4 file (this may take 30-60 seconds)...")
        final_clip.write_videofile(
            str(video_path), 
            fps=24,
            codec='libx264',
            audio_codec='aac',
            temp_audiofile=str(temp_dir / f"{video_id}_temp_audio.m4a"),
            remove_temp=True,
            verbose=False,
            logger=None
        )
        print(f"✅ Video file created: {video_path}")
        
        # Clean up temporary slide files
        for slide_path in slide_paths:
            try:
                os.remove(slide_path)
            except:
                pass
        print("✅ Cleaned up temporary slide files")
        
        # Clean up audio file
        try:
            os.remove(audio_path)
        except:
            pass
        print("✅ Cleaned up temporary audio file")
        
        print(f"🎉 Video generation complete: {video_path}")
        return str(video_path)
        
    except Exception as e:
        print(f"❌ Error during video generation: {e}")
        import traceback
        traceback.print_exc()
        raise


def generate_marketing_script(dashboard_data, timeframe, style):
    """Generate compelling Hebrew marketing script"""
    
    total_families = dashboard_data.get('total_families', 0)
    active_tutorships = dashboard_data.get('active_tutorships', 0)
    waiting_families = dashboard_data.get('waiting_families', 0)
    pending_tutors = dashboard_data.get('pending_tutors', 0)
    new_families = dashboard_data.get('new_families', 0)
    
    if style == 'מקצועי ורשמי':
        script = f"""
שלום וברוכים הבאים לסקירת מערכת חיוך של ילד.
{timeframe}, ראינו התקדמות מרשימה במערכת שלנו.
כיום, {total_families} משפחות רשומות במערכת.
{active_tutorships} חונכויות פעילות פועלות ברגע זה.
{pending_tutors} חונכים מסורים ממתינים לקבל חניכים חדשים.
{new_families} משפחות חדשות הצטרפו למערכת לאחרונה.
יחד, אנחנו יוצרים שינוי אמיתי בחייהם של ילדים ומשפחות.
תודה רבה על היותכם חלק ממסע מרגש זה.
        """
    elif style == 'ידידותי וחם':
        script = f"""
היי! שמחים לעדכן אתכם על המצב במערכת חיוך של ילד!
{timeframe} היה פשוט נהדר!
יש לנו {total_families} משפחות נפלאות במערכת.
{active_tutorships} חונכויות מדהימות פועלות עכשיו.
{pending_tutors} חונכים מצויינים מחכים לפגוש את החניכים שלהם.
ו-{new_families} משפחות חדשות הצטרפו אלינו!
אתם עושים עבודה מדהימה. המשיכו ככה!
        """
    else:  # אנרגטי ומעורר השראה
        script = f"""
וואו! בואו נראה את ההישגים המדהימים של חיוך של ילד!
{timeframe} - מה תקופה!
{total_families} משפחות! זה אדיר!
{active_tutorships} חונכויות פעילות שמשנות חיים!
{pending_tutors} חונכים נלהבים מוכנים לפעולה!
{new_families} משפחות חדשות הצטרפו למהפכה!
יחד אנחנו בלתי עצירים! בואו נמשיך להאיר חיוכים!
        """
    
    return script.strip()


def create_professional_slides(dashboard_data, timeframe):
    """Create beautiful slides with system data using Pillow with RTL support"""
    from PIL import Image, ImageDraw, ImageFont
    import os

    try:
        from bidi.algorithm import get_display
        rtl_support = True
    except ImportError:
        rtl_support = False
        print("⚠️  Warning: python-bidi not installed, Hebrew text may be mirrored")

    BASE_DIR = os.path.dirname(os.path.abspath(__file__))
    FONT_DIR = os.path.join(BASE_DIR, "fonts")

    FONTS = {
        "regular": os.path.join(FONT_DIR, "Alef-Regular.ttf"),
        "bold": os.path.join(FONT_DIR, "Alef-Bold.ttf"),
    }

    def load_font(kind: str, size: int) -> ImageFont.FreeTypeFont:
        path = FONTS[kind]
        if not os.path.exists(path):
            raise RuntimeError(f"Font missing: {path}")
        return ImageFont.truetype(path, size)

    def fix_hebrew_text(text):
        """Convert Hebrew text to display-ready format (RTL support)"""
        if rtl_support:
            return get_display(text)
        return text

    slides = []
    width, height = 1280, 720
    
    use_system_fonts = os.environ.get("USE_SYSTEM_FONTS", "false").lower() == "true"

    if use_system_fonts:
        try:
            font_large = ImageFont.truetype(
                "/System/Library/Fonts/Supplemental/Arial Unicode.ttf", 60
            )
            font_medium = ImageFont.truetype(
                "/System/Library/Fonts/Supplemental/Arial Unicode.ttf", 40
            )
            font_small = ImageFont.truetype(
                "/System/Library/Fonts/Supplemental/Arial Unicode.ttf", 32
            )
        except Exception:
            try:
                font_large = ImageFont.truetype('/System/Library/Fonts/Supplemental/Arial.ttf', 60)
                font_medium = ImageFont.truetype('/System/Library/Fonts/Supplemental/Arial.ttf', 40)
                font_small = ImageFont.truetype('/System/Library/Fonts/Supplemental/Arial.ttf', 32)
            except Exception:
                # Fallback to bundled fonts if system fonts fail
                font_large = load_font("bold", 60)
                font_medium = load_font("regular", 40)
                font_small = load_font("regular", 32)
    else:
        # ✅ Production-safe, deterministic path
        font_large = load_font("bold", 60)
        font_medium = load_font("regular", 40)
        font_small = load_font("regular", 32)

    # Slide 1: Title Slide
    img1 = Image.new('RGB', (width, height), color='#667EEA')
    draw1 = ImageDraw.Draw(img1)
    draw1.text((640, 280), 'ChildSmile', font=font_large, fill='white', anchor='mm')
    draw1.text((640, 380), fix_hebrew_text(f'סקירת מערכת - {timeframe}'), font=font_medium, fill='#E0E7FF', anchor='mm')
    slides.append(img1)
    
    # Slide 2: KPIs
    img2 = Image.new('RGB', (width, height), color='#1E293B')
    draw2 = ImageDraw.Draw(img2)
    draw2.text((640, 80), fix_hebrew_text('מדדים עיקריים'), font=font_medium, fill='#667EEA', anchor='mm')
    
    kpis = [
        (fix_hebrew_text(f"סה\"כ משפחות: {dashboard_data.get('total_families', 0)}"), '#3B82F6'),
        (fix_hebrew_text(f"חונכויות פעילות: {dashboard_data.get('active_tutorships', 0)}"), '#10B981'),
        (fix_hebrew_text(f"משפחות ממתינות: {dashboard_data.get('waiting_families', 0)}"), '#F59E0B'),
        (fix_hebrew_text(f"חונכים ממתינים: {dashboard_data.get('pending_tutors', 0)}"), '#8B5CF6'),
    ]
    
    y_pos = 200
    for kpi_text, color in kpis:
        draw2.text((640, y_pos), kpi_text, font=font_small, fill=color, anchor='mm')
        y_pos += 100
    
    slides.append(img2)
    
    # Slide 3: New Families
    img3 = Image.new('RGB', (width, height), color='#0F172A')
    draw3 = ImageDraw.Draw(img3)
    draw3.text((640, 250), fix_hebrew_text('משפחות חדשות'), font=font_medium, fill='#667EEA', anchor='mm')
    new_fam = dashboard_data.get('new_families', 0)
    draw3.text((640, 380), f'{new_fam}', font=font_large, fill='#10B981', anchor='mm')
    draw3.text((640, 480), fix_hebrew_text('הצטרפו לאחרונה'), font=font_small, fill='white', anchor='mm')
    slides.append(img3)
    
    # Slide 4: Thank You
    img4 = Image.new('RGB', (width, height), color='#667EEA')
    draw4 = ImageDraw.Draw(img4)
    draw4.text((640, 320), fix_hebrew_text('תודה רבה!'), font=font_large, fill='white', anchor='mm')
    draw4.text((640, 420), fix_hebrew_text('יחד אנחנו עושים את ההבדל'), font=font_small, fill='#E0E7FF', anchor='mm')
    slides.append(img4)
    
    return slides


def generate_ppt_slide(ppt_id, dashboard_data):
    """
    Generate PowerPoint presentation
    Using python-pptx library
    """
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except ImportError:
        # If python-pptx not installed, create placeholder
        temp_dir = Path(tempfile.gettempdir()) / 'childsmile_ppts'
        temp_dir.mkdir(exist_ok=True)
        ppt_path = temp_dir / f"{ppt_id}.pptx"
        
        # Create empty file as placeholder
        with open(ppt_path, 'wb') as f:
            f.write(b'Placeholder PPT - Install python-pptx for real generation')
        
        return str(ppt_path)
    
    # Create temp directory
    temp_dir = Path(tempfile.gettempdir()) / 'childsmile_ppts'
    temp_dir.mkdir(exist_ok=True)
    
    ppt_path = temp_dir / f"{ppt_id}.pptx"
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Wide format
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (simulated with solid color)
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(102, 126, 234)
    
    # Add title
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "לוח בקרה - ChildSmile"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(4), Inches(11.33), Inches(0.5)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "סקירת מצב מערכת"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # KPIs Slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box2 = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(12), Inches(0.5)
    )
    title_frame2 = title_box2.text_frame
    title_frame2.text = "מדדים עיקריים"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(32)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(102, 126, 234)
    
    # Add table
    rows, cols = 6, 2
    left = Inches(2)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    
    table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column widths
    table.columns[0].width = Inches(6)
    table.columns[1].width = Inches(3)
    
    # Fill header
    table.cell(0, 0).text = "מדד"
    table.cell(0, 1).text = "ערך"
    
    # Fill data
    kpi_data = [
        ["סה\"כ משפחות", str(dashboard_data['total_families'])],
        ["משפחות הממתינות", str(dashboard_data['waiting_families'])],
        ["חונכויות פעילות", str(dashboard_data['active_tutorships'])],
        ["חונכים הממתינים", str(dashboard_data['pending_tutors'])],
        ["אנשי הצוות", str(dashboard_data['staff_count'])]
    ]
    
    for i, (metric, value) in enumerate(kpi_data, start=1):
        table.cell(i, 0).text = metric
        table.cell(i, 1).text = value
        
        # Style cells
        for j in range(2):
            cell = table.cell(i, j)
            cell.text_frame.paragraphs[0].font.size = Pt(18)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
    
    # Style header
    for j in range(2):
        cell = table.cell(0, j)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(20)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Save presentation
    prs.save(str(ppt_path))
    
    return str(ppt_path)


def cleanup_temp_files(file_path, delay_minutes=60):
    """
    Schedule cleanup of temporary files
    Files are deleted after specified delay
    """
    
    def delete_file():
        time.sleep(delay_minutes * 60)
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
                print(f"Cleaned up temp file: {file_path}")
        except Exception as e:
            print(f"Error cleaning up {file_path}: {e}")
    
    # Run cleanup in background thread
    cleanup_thread = threading.Thread(target=delete_file, daemon=True)
    cleanup_thread.start()
